"""
知识库管理界面 — 查看/编辑知识结构、导入PDF/PPT并AI辅助分类
"""
from __future__ import annotations

import json
import os
import re
from pathlib import Path

from PySide6.QtCore import Qt
from PySide6.QtGui import QAction
from PySide6.QtWidgets import (
    QComboBox,
    QFileDialog,
    QHBoxLayout,
    QLabel,
    QLineEdit,
    QMenu,
    QMessageBox,
    QPushButton,
    QSplitter,
    QTextEdit,
    QTreeWidget,
    QTreeWidgetItem,
    QVBoxLayout,
    QWidget,
)

from raggg.config import Settings
from raggg.i18n import get_text
from raggg.theme import get_colors

COLORS = get_colors()


class KnowledgeManager(QWidget):

    def __init__(self, settings: Settings, parent: QWidget | None = None) -> None:
        super().__init__(parent)
        self.settings = settings
        self._kb_root = Path(__file__).resolve().parents[3] / "knowledge_base"
        self._current_file: Path | None = None
        self._pending_import_path: str = ""
        self._pending_import_type: str = ""
        self._build_ui()
        self._load_tree()

    # ─── UI ─────────────────────────────────────
    def _build_ui(self) -> None:
        layout = QVBoxLayout(self)
        layout.setSpacing(10)

        desc = QLabel(get_text("kbm_desc"))
        desc.setStyleSheet(f"color:{COLORS['muted']};font-size:12px;")
        desc.setWordWrap(True)
        layout.addWidget(desc)

        splitter = QSplitter(Qt.Horizontal)

        # 左侧：树 + 右键菜单
        left = QWidget()
        left_layout = QVBoxLayout(left)
        left_layout.setContentsMargins(0, 0, 0, 0)
        left_layout.setSpacing(6)

        tree_label = QLabel(get_text("kbm_tree_title"))
        tree_label.setStyleSheet(f"color:{COLORS['text']};font-weight:700;")
        left_layout.addWidget(tree_label)

        self.tree = QTreeWidget()
        self.tree.setHeaderHidden(True)
        self.tree.setContextMenuPolicy(Qt.CustomContextMenu)
        self.tree.itemClicked.connect(self._on_tree_select)
        self.tree.customContextMenuRequested.connect(self._on_tree_context_menu)
        left_layout.addWidget(self.tree)

        splitter.addWidget(left)

        # 右侧：编辑区
        right = QWidget()
        right_layout = QVBoxLayout(right)
        right_layout.setContentsMargins(0, 0, 0, 0)
        right_layout.setSpacing(6)

        info_bar = QHBoxLayout()
        self.file_label = QLabel(get_text("kbm_select_hint"))
        self.file_label.setStyleSheet(f"color:{COLORS['muted']};font-size:11px;")
        info_bar.addWidget(self.file_label, stretch=1)

        save_btn = QPushButton(get_text("kbm_save"))
        save_btn.clicked.connect(self._save_current)
        info_bar.addWidget(save_btn)
        right_layout.addLayout(info_bar)

        self.editor = QTextEdit()
        self.editor.setPlaceholderText(get_text("kbm_editor_placeholder"))
        right_layout.addWidget(self.editor)

        splitter.addWidget(right)
        splitter.setSizes([280, 480])
        layout.addWidget(splitter, stretch=1)

        # ── 导入区 ──
        import_section = QWidget()
        import_layout = QVBoxLayout(import_section)
        import_layout.setContentsMargins(0, 0, 0, 0)
        import_layout.setSpacing(6)

        import_label = QLabel(get_text("kbm_import_title"))
        import_label.setStyleSheet(f"color:{COLORS['text']};font-weight:700;")
        import_layout.addWidget(import_label)

        # Row 1: 选择文件
        row1 = QHBoxLayout(); row1.setSpacing(8)
        self.file_select_btn = QPushButton(get_text("kbm_select_file"))
        self.file_select_btn.clicked.connect(self._select_file)
        row1.addWidget(self.file_select_btn)
        self.file_status = QLabel(get_text("kbm_no_file"))
        self.file_status.setStyleSheet(f"color:{COLORS['muted']};font-size:11px;")
        row1.addWidget(self.file_status, stretch=1)
        import_layout.addLayout(row1)

        # Row 2: 分类 + 子目录 + 重要度
        row2 = QHBoxLayout(); row2.setSpacing(8)
        row2.addWidget(QLabel(get_text("kbm_import_target") + ":"))
        self.category_combo = QComboBox()
        self._populate_categories()
        self.category_combo.currentIndexChanged.connect(self._on_category_changed)
        row2.addWidget(self.category_combo)

        row2.addWidget(QLabel(get_text("kbm_import_subdir") + ":"))
        self.subdir_combo = QComboBox()
        self._populate_subdirs()
        row2.addWidget(self.subdir_combo)

        row2.addWidget(QLabel(get_text("kbm_priority_label") + ":"))
        self.priority_combo = QComboBox()
        for i in range(1, 6):
            self.priority_combo.addItem(get_text(f"kbm_priority_{i}"), i)
        self.priority_combo.setCurrentIndex(2)
        row2.addWidget(self.priority_combo)
        import_layout.addLayout(row2)

        # Row 3: 备注 + 确认导入
        row3 = QHBoxLayout(); row3.setSpacing(8)
        row3.addWidget(QLabel(get_text("kbm_import_desc_label") + " *:"))
        self.desc_edit = QLineEdit()
        self.desc_edit.setPlaceholderText(get_text("kbm_import_desc_placeholder"))
        row3.addWidget(self.desc_edit, stretch=1)

        self.confirm_btn = QPushButton(get_text("kbm_confirm_import_btn"))
        self.confirm_btn.clicked.connect(self._confirm_import)
        row3.addWidget(self.confirm_btn)
        import_layout.addLayout(row3)

        layout.addWidget(import_section)

    # ─── Tree ──────────────────────────────────
    _CAT_I18N = {
        "01_team_tutorials": "kbm_cat_01", "02_software_manual": "kbm_cat_02",
        "03_examples": "kbm_cat_03", "04_error_cases": "kbm_cat_04",
        "05_reference": "kbm_cat_05", "06_theory_notes": "kbm_cat_06",
        "tutorials": "kbm_sub_tutorials", "Circuit": "kbm_example_circuit",
        "EM": "kbm_example_em", "Mech": "kbm_example_mech",
        "Multi-Physics": "kbm_example_multi", "Thermal": "kbm_example_thermal",
    }

    def _load_tree(self) -> None:
        self.tree.clear()
        if self._kb_root.exists():
            self._add_tree_items(self.tree.invisibleRootItem(), self._kb_root)

    def _add_tree_items(self, parent: QTreeWidgetItem, path: Path) -> None:
        for entry in sorted(path.iterdir()):
            if entry.name.startswith("."):
                continue
            if entry.is_dir():
                display_name = get_text(self._CAT_I18N.get(entry.name, "")) or entry.name
                item = QTreeWidgetItem(parent, [display_name])
                item.setData(0, Qt.UserRole, str(entry))
                md_count = sum(1 for _ in entry.rglob("*.md"))
                item.setToolTip(0, f"{md_count} documents")
                self._add_tree_items(item, entry)
            elif entry.suffix == ".md":
                item = QTreeWidgetItem(parent, [entry.stem])
                item.setData(0, Qt.UserRole, str(entry))
                item.setToolTip(0, str(entry))
        if parent == self.tree.invisibleRootItem():
            self.tree.expandAll()

    def _on_tree_select(self, item: QTreeWidgetItem, _column: int) -> None:
        filepath = item.data(0, Qt.UserRole)
        if not filepath:
            return
        path = Path(filepath)
        if path.is_file() and path.suffix == ".md":
            self._load_file(path)

    def _load_file(self, path: Path) -> None:
        self._current_file = path
        self.file_label.setText(str(path.relative_to(self._kb_root.parent)))
        try:
            text = path.read_text(encoding="utf-8")
        except Exception:
            text = path.read_text(encoding="utf-8", errors="ignore")
        self.editor.setPlainText(text)

    def _save_current(self) -> None:
        if not self._current_file:
            QMessageBox.information(self, get_text("kbm_save"), get_text("kbm_no_file_selected"))
            return
        self._current_file.write_text(self.editor.toPlainText(), encoding="utf-8")
        QMessageBox.information(self, get_text("kbm_save"), get_text("kbm_save_ok"))

    # ─── 右键菜单：删除文件 ─────────────────────
    def _on_tree_context_menu(self, pos) -> None:
        item = self.tree.itemAt(pos)
        if not item:
            return
        filepath = item.data(0, Qt.UserRole)
        if not filepath:
            return
        path = Path(filepath)
        if not path.is_file():
            return

        menu = QMenu(self)
        del_action = QAction(get_text("kbm_delete_file"), self)
        del_action.triggered.connect(lambda: self._delete_file(path, item))
        menu.addAction(del_action)
        menu.exec(self.tree.viewport().mapToGlobal(pos))

    def _delete_file(self, path: Path, item: QTreeWidgetItem) -> None:
        name = path.name
        reply = QMessageBox.question(
            self,
            get_text("kbm_delete_confirm_title"),
            get_text("kbm_delete_confirm_msg").replace("{name}", name),
            QMessageBox.Yes | QMessageBox.No,
        )
        if reply != QMessageBox.Yes:
            return
        try:
            path.unlink()
            parent = item.parent() or self.tree.invisibleRootItem()
            parent.removeChild(item)
            if self._current_file and str(self._current_file) == str(path):
                self._current_file = None
                self.editor.clear()
                self.file_label.setText(get_text("kbm_select_hint"))
            # 移除空目录
            parent_path = Path(path.parent)
            if parent_path != self._kb_root and not any(parent_path.iterdir()):
                parent_path.rmdir()
                self._load_tree()
            # 重建索引
            from raggg.pipeline.builder import build_knowledge_base
            build_knowledge_base(self.settings)
        except Exception as e:
            QMessageBox.warning(self, get_text("kbm_delete_error"), str(e))

    # ─── Import ────────────────────────────────
    def _populate_categories(self) -> None:
        cats = [
            ("01_team_tutorials", get_text("kbm_cat_01")),
            ("02_software_manual", get_text("kbm_cat_02")),
            ("03_examples", get_text("kbm_cat_03")),
            ("04_error_cases", get_text("kbm_cat_04")),
            ("05_reference", get_text("kbm_cat_05")),
            ("06_theory_notes", get_text("kbm_cat_06")),
        ]
        self.category_combo.clear()
        for cat_id, cat_label in cats:
            self.category_combo.addItem(cat_label, cat_id)

    def _on_category_changed(self, _index: int) -> None:
        self._populate_subdirs()

    def _populate_subdirs(self) -> None:
        cat_id = self.category_combo.currentData()
        cat_path = self._kb_root / cat_id
        self.subdir_combo.clear()
        self.subdir_combo.addItem(get_text("kbm_subdir_auto"), "__AUTO__")
        self.subdir_combo.addItem(get_text("kbm_subdir_root"), "")
        if cat_path.exists():
            self._add_subdir_options(self.subdir_combo, cat_path, "")
        self.subdir_combo.setCurrentIndex(0)

    def _add_subdir_options(self, combo: QComboBox, base: Path, prefix: str, depth: int = 0) -> None:
        """递归添加所有深度的子目录，用缩进体现层级"""
        indent = "  " * depth + ("" if depth == 0 else "├ ")
        for entry in sorted(base.iterdir()):
            if entry.is_dir() and not entry.name.startswith("."):
                display = self._CAT_I18N.get(entry.name, entry.name)
                rel_path = (prefix + "/" + entry.name).lstrip("/")
                combo.addItem(f"{indent}{display}", rel_path)
                self._add_subdir_options(combo, entry, rel_path, depth + 1)

    # ── 文件选择 ────────────────────────────────
    def _select_file(self) -> None:
        path, _ = QFileDialog.getOpenFileName(
            self, get_text("kbm_select_file_title"), "",
            "Documents (*.pdf *.pptx)"
        )
        if not path:
            return
        self._pending_import_path = path
        self._pending_import_type = "pdf" if path.lower().endswith(".pdf") else "pptx"
        self.file_status.setText(f"📎 {Path(path).name}")
        self.file_status.setStyleSheet(f"color:{COLORS['accent']};font-size:11px;")

    # ── 确认导入（强制备注） ─────────────────────
    def _confirm_import(self) -> None:
        if not self._pending_import_path:
            QMessageBox.information(self, get_text("kbm_no_file"), get_text("kbm_no_file_msg"))
            return
        user_desc = self.desc_edit.text().strip()
        if not user_desc:
            QMessageBox.information(self, get_text("kbm_notes_required"), get_text("kbm_notes_required_msg"))
            return
        self._import_file(self._pending_import_path, self._pending_import_type)
        self._pending_import_path = ""
        self._pending_import_type = ""
        self.file_status.setText(get_text("kbm_no_file"))
        self.file_status.setStyleSheet(f"color:{COLORS['muted']};font-size:11px;")

    def _import_file(self, filepath: str, filetype: str) -> None:
        filepath = Path(filepath)

        # 1. 提取文本
        if filetype == "pdf":
            try:
                from raggg.pipeline.ingestion import _extract_pdf_text
                text = _extract_pdf_text(filepath)
            except ImportError:
                text = f"[PDF: {filepath.name}]\n\n(需要安装 pypdf)"
        else:
            text = self._extract_pptx_text(filepath)

        if not text.strip():
            QMessageBox.warning(self, get_text("kbm_import_failed"), get_text("kbm_import_empty"))
            return

        # 2. 生成 Markdown
        markdown = self._build_markdown(filepath, text)

        # 3. 目标路径（大分类 + AI/手动子目录，支持多级深度）
        cat_id = self.category_combo.currentData()
        base_dir = self._kb_root / cat_id
        user_desc = self.desc_edit.text().strip()
        subdir_choice = self.subdir_combo.currentData()

        if subdir_choice == "__AUTO__":
            all_subdirs = self._list_all_subdirs(base_dir)
            chosen_subdir = self._ai_suggest_subdir(text, filepath.stem, user_desc, all_subdirs)
            target_dir = base_dir / chosen_subdir if chosen_subdir else base_dir
        elif subdir_choice:
            target_dir = base_dir / subdir_choice
        else:
            target_dir = base_dir

        # 4. AI 分析
        suggestion = self._ai_analyze(text, filepath.stem, user_desc)

        # 5. 预览并确认
        preview = (
            f"{get_text('kbm_preview_file')}: {filepath.name}\n"
            f"{get_text('kbm_preview_target')}: {target_dir.relative_to(self._kb_root.parent)}\n"
            f"{get_text('kbm_preview_suggestion')}:\n{suggestion}\n\n"
            f"---\n{markdown[:500]}...\n"
        )
        reply = QMessageBox.question(
            self, get_text("kbm_confirm_import"), preview,
            QMessageBox.Yes | QMessageBox.No,
        )
        if reply != QMessageBox.Yes:
            return

        # 6. 写入
        target_dir.mkdir(parents=True, exist_ok=True)
        safe_name = re.sub(r"[^\w\-.]", "_", filepath.stem)
        out_path = target_dir / f"{safe_name}.md"
        out_path.write_text(markdown, encoding="utf-8")
        self._load_tree()

        # 7. 重建索引
        from raggg.pipeline.builder import build_knowledge_base
        try:
            report = build_knowledge_base(self.settings)
            msg = f"{get_text('kbm_import_ok_msg')}: {out_path.relative_to(self._kb_root.parent)}\nChunks: {report.chunk_count}"
        except Exception:
            msg = f"{get_text('kbm_import_ok_msg')}: {out_path.relative_to(self._kb_root.parent)}\n(索引重建失败)"
        QMessageBox.information(self, get_text("kbm_import_done"), msg)

    def _build_markdown(self, filepath: Path, text: str) -> str:
        cat_id = self.category_combo.currentData()
        priority = self.priority_combo.currentData()
        return (
            "---\n"
            f"title: \"{filepath.stem}\"\n"
            f"content_kind: \"imported\"\n"
            f"source_file: \"{filepath.name}\"\n"
            f"imported_at: \"2026-07-13\"\n"
            f"category: \"{cat_id}\"\n"
            f"priority: {priority}\n"
            "---\n\n"
            f"# {filepath.stem}\n\n{text}\n"
        )

    @staticmethod
    def _list_all_subdirs(base_dir: Path) -> list[str]:
        """递归列出所有深度子目录的相对路径"""
        subdirs = []
        if not base_dir.exists():
            return subdirs
        for entry in sorted(base_dir.iterdir()):
            if entry.is_dir() and not entry.name.startswith("."):
                subdirs.append(entry.name)
                for child in entry.rglob("*"):
                    if child.is_dir() and not child.name.startswith("."):
                        rel = str(child.relative_to(base_dir)).replace("\\", "/")
                        subdirs.append(rel)
        return subdirs

    def _ai_suggest_subdir(
        self, text: str, filename: str, user_desc: str, subdirs: list[str]
    ) -> str:
        """让 AI 从所有深度子目录中推荐最佳匹配"""
        if not subdirs:
            return ""
        prompt = (
            f"分析文档内容，从候选目录中选择最匹配的一个（支持多级子目录如 Modeling/Stimulate）。\n"
            f"文档名：{filename}\n"
            f"候选目录：{', '.join(subdirs)}\n"
            f"内容：{text[:1500]}\n"
        )
        if user_desc:
            prompt += f"管理员备注：{user_desc}\n"
        prompt += "只回答目录路径，不要解释。不确定则回答 ROOT。"

        if self.settings.llm_api_key:
            try:
                from raggg.generation.llm_client import OpenAICompatibleClient
                client = OpenAICompatibleClient(
                    self.settings.llm_base_url, self.settings.llm_api_key, self.settings.llm_model,
                )
                response = client.complete(prompt).strip()
                if response in subdirs:
                    return response
            except Exception:
                pass

        # 关键词回退（支持多级匹配）
        for sub in sorted(subdirs, key=lambda s: -len(s)):  # 深的优先
            lower = sub.lower()
            keywords = [sub.split("/")[-1].lower()]
            if "em" in lower or "project" in lower:
                keywords += ["端口", "激励", "边界", "远场", "pml", "s参数"]
            if "modeling" in lower:
                keywords += ["建模", "几何", "曲线", "拉伸"]
            if "stimulate" in lower:
                keywords += ["端口", "激励", "lumped", "wave", "plane"]
            if "mesh" in lower:
                keywords += ["网格", "mesh"]
            if "antenna" in lower:
                keywords += ["天线"]
            if "filter" in lower:
                keywords += ["滤波器"]
            if "design" in lower:
                keywords += ["domain", "频率", "单位", "求解器"]
            for kw in keywords:
                if kw.lower() in text.lower() or kw.lower() in filename.lower():
                    return sub
        return ""

    def _ai_analyze(self, text: str, filename: str, user_desc: str) -> str:
        prompt = (
            "你是知识库管理助手。分析文档内容，给出：\n"
            "1. 建议分类\n2. 3-5 关键词\n3. 一句话摘要\n\n"
        )
        if user_desc:
            prompt += f"管理员备注：{user_desc}\n\n"
        prompt += f"文档名：{filename}\n内容片段：{text[:1500]}\n"
        prompt += "格式：分类: xxx | 关键词: a,b,c | 摘要: xxx"

        if self.settings.llm_api_key:
            try:
                from raggg.generation.llm_client import OpenAICompatibleClient
                client = OpenAICompatibleClient(
                    self.settings.llm_base_url, self.settings.llm_api_key, self.settings.llm_model,
                )
                return client.complete(prompt).strip()
            except Exception:
                pass

        keywords_map = {
            "端口": "02_software_manual", "激励": "02", "边界": "02", "求解器": "02",
            "案例": "03_examples", "仿真": "03", "错误": "04_error_cases",
            "教程": "01_team_tutorials", "FAQ": "01",
            "材料": "05_reference", "理论": "06_theory_notes",
        }
        for kw, cat in keywords_map.items():
            if kw in text or kw in filename:
                return f"分类: {cat} | 关键词: {kw} | 摘要: (本地匹配)"
        return "分类: 02_software_manual | 关键词: 待补充 | 摘要: (本地匹配)"

    @staticmethod
    def _extract_pptx_text(filepath: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            return f"[PPT: {filepath.name}]\n\n(需要安装 python-pptx)"
        prs = Presentation(str(filepath))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            lines = [f"## Slide {i}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lines.append(t)
                if shape.has_table:
                    rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows]
                    if rows:
                        lines.append("\n| " + " |\n| ".join(rows) + " |")
            if len(lines) > 1:
                slides.append("\n".join(lines))
        return "\n\n".join(slides) if slides else f"[PPT: {filepath.name} - 无可提取文本]"

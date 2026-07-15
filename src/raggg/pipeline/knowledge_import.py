from __future__ import annotations

from datetime import date
from pathlib import Path

from raggg.config import Settings


def build_import_markdown(
    filepath: Path,
    text: str,
    category: str,
    priority: int,
) -> str:
    return (
        "---\n"
        f'title: "{filepath.stem}"\n'
        'content_kind: "imported"\n'
        f'source_file: "{filepath.name}"\n'
        f'imported_at: "{date.today().isoformat()}"\n'
        f'category: "{category}"\n'
        f"priority: {priority}\n"
        "---\n\n"
        f"# {filepath.stem}\n\n{text}\n"
    )


def list_all_subdirs(base_dir: Path) -> list[str]:
    if not base_dir.exists():
        return []
    return [
        path.relative_to(base_dir).as_posix()
        for path in sorted(base_dir.rglob("*"))
        if path.is_dir() and not any(part.startswith(".") for part in path.relative_to(base_dir).parts)
    ]


def suggest_subdir(
    settings: Settings,
    text: str,
    filename: str,
    user_desc: str,
    subdirs: list[str],
) -> str:
    if not subdirs:
        return ""
    prompt = (
        "分析文档内容，从候选目录中选择最匹配的一个（支持多级子目录如 Modeling/Stimulate）。\n"
        f"文档名：{filename}\n"
        f"候选目录：{', '.join(subdirs)}\n"
        f"内容：{text[:1500]}\n"
    )
    if user_desc:
        prompt += f"管理员备注：{user_desc}\n"
    prompt += "只回答目录路径，不要解释。不确定则回答 ROOT。"

    response = _complete_if_configured(settings, prompt)
    if response in subdirs:
        return response

    searchable = f"{filename}\n{text}".lower()
    for subdir in sorted(subdirs, key=lambda value: (-value.count("/"), -len(value))):
        lower = subdir.lower()
        keywords = [subdir.split("/")[-1].lower()]
        if "em" in lower or "project" in lower:
            keywords += ["端口", "激励", "边界", "远场", "pml", "s参数"]
        if "modeling" in lower:
            keywords += ["建模", "几何", "曲线", "拉伸"]
        if "stimulate" in lower:
            keywords += ["端口", "激励", "lumped", "wave", "plane"]
        if "mesh" in lower:
            keywords += ["网格", "mesh"]
        if "antenna" in lower:
            keywords.append("天线")
        if "filter" in lower:
            keywords.append("滤波器")
        if "design" in lower:
            keywords += ["domain", "频率", "单位", "求解器"]
        if any(keyword in searchable for keyword in keywords):
            return subdir
    return ""


def analyze_document(
    settings: Settings,
    text: str,
    filename: str,
    user_desc: str,
) -> str:
    prompt = (
        "你是知识库管理助手。分析文档内容，给出：\n"
        "1. 建议分类\n2. 3-5 关键词\n3. 一句话摘要\n\n"
    )
    if user_desc:
        prompt += f"管理员备注：{user_desc}\n\n"
    prompt += f"文档名：{filename}\n内容片段：{text[:1500]}\n"
    prompt += "格式：分类: xxx | 关键词: a,b,c | 摘要: xxx"

    response = _complete_if_configured(settings, prompt)
    if response:
        return response

    keywords_map = {
        "端口": "02_software_manual",
        "激励": "02_software_manual",
        "边界": "02_software_manual",
        "求解器": "02_software_manual",
        "案例": "03_examples",
        "仿真": "03_examples",
        "错误": "04_error_cases",
        "教程": "01_team_tutorials",
        "FAQ": "01_team_tutorials",
        "材料": "05_reference",
        "理论": "06_theory_notes",
    }
    for keyword, category in keywords_map.items():
        if keyword in text or keyword in filename:
            return f"分类: {category} | 关键词: {keyword} | 摘要: (本地匹配)"
    return "分类: 02_software_manual | 关键词: 待补充 | 摘要: (本地匹配)"


def extract_pptx_text(filepath: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return f"[PPT: {filepath.name}]\n\n(需要安装 python-pptx)"
    presentation = Presentation(str(filepath))
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, 1):
        lines = [f"## Slide {index}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.extend(
                    paragraph.text.strip()
                    for paragraph in shape.text_frame.paragraphs
                    if paragraph.text.strip()
                )
            if shape.has_table:
                rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows]
                if rows:
                    lines.append("\n| " + " |\n| ".join(rows) + " |")
        if len(lines) > 1:
            slides.append("\n".join(lines))
    return "\n\n".join(slides) if slides else f"[PPT: {filepath.name} - 无可提取文本]"


def _complete_if_configured(settings: Settings, prompt: str) -> str:
    if not settings.llm_api_key:
        return ""
    try:
        from raggg.generation.llm_client import OpenAICompatibleClient

        client = OpenAICompatibleClient(
            settings.llm_base_url,
            settings.llm_api_key,
            settings.llm_model,
        )
        return client.complete(prompt).strip()
    except Exception:
        return ""
